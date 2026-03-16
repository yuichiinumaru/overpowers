from pptx import Presentation
import json
import os
import argparse

def extract_pptx(file_path, output_dir):
    """
    Extract all content from a PowerPoint file.
    Returns a JSON structure with slides, text, and images.
    """
    prs = Presentation(file_path)
    slides_data = []

    # Create assets directory
    assets_dir = os.path.join(output_dir, 'assets')
    os.makedirs(assets_dir, exist_ok=True)

    for slide_num, slide in enumerate(prs.slides):
        slide_data = {
            'number': slide_num + 1,
            'title': '',
            'content': [],
            'images': [],
            'notes': ''
        }

        for shape in slide.shapes:
            # Extract title
            if shape.has_text_frame:
                if shape == slide.shapes.title:
                    slide_data['title'] = shape.text
                else:
                    slide_data['content'].append({
                        'type': 'text',
                        'content': shape.text
                    })

            # Extract images
            if shape.shape_type == 13:  # Picture
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                image_name = f"slide{slide_num + 1}_img{len(slide_data['images']) + 1}.{image_ext}"
                image_path = os.path.join(assets_dir, image_name)

                with open(image_path, 'wb') as f:
                    f.write(image_bytes)

                slide_data['images'].append({
                    'path': f"assets/{image_name}",
                    'width': shape.width,
                    'height': shape.height
                })

        # Extract notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            slide_data['notes'] = notes_frame.text

        slides_data.append(slide_data)

    return slides_data

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Extract content from PPTX file")
    parser.add_argument("file_path", help="Path to PPTX file")
    parser.add_argument("output_dir", help="Directory to save assets and JSON")
    args = parser.parse_args()
    
    print(f"Extracting {args.file_path} to {args.output_dir}...")
    try:
        data = extract_pptx(args.file_path, args.output_dir)
        json_path = os.path.join(args.output_dir, "extracted_content.json")
        with open(json_path, 'w') as f:
            json.dump(data, f, indent=2)
        print(f"Extraction complete. JSON saved to {json_path}")
    except Exception as e:
        print(f"Error extracting PPTX: {e}. Ensure python-pptx is installed.")
