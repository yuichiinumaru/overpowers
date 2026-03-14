#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Markdown 到 PPT 转换器
支持将 Markdown 转换为 WPS 演示文稿
"""

import re
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


class MarkdownToPPTConverter:
    """Markdown 到 PPT 转换器"""
    
    def __init__(self):
        self.prs = None
        self.default_font = '微软雅黑'
        self.title_font_size = Pt(32)
        self.subtitle_font_size = Pt(18)
        self.content_font_size = Pt(18)
        self.code_font_size = Pt(14)
        
    def convert(self, md_content: str, output_path: str, title: str = None):
        """
        将 Markdown 内容转换为 PPT 文件
        
        Args:
            md_content: Markdown 文本内容
            output_path: 输出 PPT 文件路径
            title: 演示文稿标题（可选）
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # 解析 Markdown 并创建幻灯片
        self._parse_markdown(md_content, title)
        
        # 保存文件
        self.prs.save(output_path)
        return output_path
    
    def convert_file(self, md_file: str, output_path: str = None, title: str = None):
        """
        将 Markdown 文件转换为 PPT 文件
        
        Args:
            md_file: Markdown 文件路径
            output_path: 输出 PPT 文件路径（可选）
            title: 演示文稿标题（可选）
        """
        # 读取 Markdown 文件
        with open(md_file, 'r', encoding='utf-8') as f:
            md_content = f.read()
        
        # 确定输出路径
        if output_path is None:
            base_name = os.path.splitext(md_file)[0]
            output_path = base_name + '.pptx'
        
        return self.convert(md_content, output_path, title)
    
    def _parse_markdown(self, md_content: str, default_title: str = None):
        """解析 Markdown 内容"""
        lines = md_content.split('\n')
        i = 0
        slide_count = 0
        
        while i < len(lines):
            line = lines[i]
            stripped = line.strip()
            
            # 空行
            if not stripped:
                i += 1
                continue
            
            # H1 - 标题幻灯片
            if stripped.startswith('# ') and not stripped.startswith('## '):
                title = stripped[2:]
                if slide_count == 0 and default_title:
                    title = default_title
                i = self._create_title_slide(title, lines, i + 1)
                slide_count += 1
                continue
            
            # H2 - 内容幻灯片
            if stripped.startswith('## '):
                title = stripped[3:]
                i = self._create_content_slide(title, lines, i + 1)
                slide_count += 1
                continue
            
            # H3-H6 - 子标题幻灯片
            if re.match(r'^#{3,6}\s+', stripped):
                title = re.sub(r'^#{3,6}\s+', '', stripped)
                i = self._create_subtitle_slide(title, lines, i + 1)
                slide_count += 1
                continue
            
            # 表格 - 表格幻灯片
            if '|' in stripped and i + 1 < len(lines) and '|' in lines[i + 1] and '---' in lines[i + 1]:
                i = self._create_table_slide(lines, i)
                slide_count += 1
                continue
            
            # 代码块 - 代码幻灯片
            if stripped.startswith('```'):
                i = self._create_code_slide(lines, i)
                slide_count += 1
                continue
            
            # 图片 - 图片幻灯片
            if stripped.startswith('!['):
                i = self._create_image_slide(stripped, lines, i)
                slide_count += 1
                continue
            
            # 普通段落
            if stripped and not stripped.startswith(('#', '- ', '* ', '+ ', '>')):
                i = self._create_paragraph_slide(stripped, lines, i + 1)
                slide_count += 1
                continue
            
            i += 1
    
    def _create_title_slide(self, title: str, lines: list, start_idx: int) -> int:
        """创建标题幻灯片"""
        # 使用空白布局
        blank_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = self.title_font_size
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.name = self.default_font
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # 查找副标题（下一段落）
        i = start_idx
        while i < len(lines):
            line = lines[i].strip()
            if line and not line.startswith('#'):
                # 添加副标题
                subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
                subtitle_frame = subtitle_box.text_frame
                subtitle_frame.text = line
                subtitle_frame.paragraphs[0].font.size = self.subtitle_font_size
                subtitle_frame.paragraphs[0].font.name = self.default_font
                subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                return i + 1
            elif line.startswith('#'):
                return i
            i += 1
        
        return i
    
    def _create_content_slide(self, title: str, lines: list, start_idx: int) -> int:
        """创建内容幻灯片"""
        # 使用空白布局
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.name = self.default_font
        
        # 收集内容
        i = start_idx
        content_lines = []
        
        while i < len(lines):
            line = lines[i]
            stripped = line.strip()
            
            # 遇到新标题结束
            if stripped.startswith('#'):
                break
            
            # 列表项
            if stripped.startswith(('- ', '* ', '+ ')):
                content = stripped[2:]
                content_lines.append(f"• {content}")
                i += 1
            elif re.match(r'^\d+\.', stripped):
                content = re.sub(r'^\d+\.\s*', '', stripped)
                content_lines.append(f"• {content}")
                i += 1
            elif stripped and not stripped.startswith('```') and not stripped.startswith('!['):
                content_lines.append(stripped)
                i += 1
            else:
                break
        
        # 添加内容
        if content_lines:
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            for idx, line in enumerate(content_lines):
                if idx == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()
                
                p.text = line
                p.font.size = self.content_font_size
                p.font.name = self.default_font
                p.space_after = Pt(12)
        
        return i
    
    def _create_subtitle_slide(self, title: str, lines: list, start_idx: int) -> int:
        """创建子标题幻灯片"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 添加子标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.name = self.default_font
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return start_idx
    
    def _create_table_slide(self, lines: list, start_idx: int) -> int:
        """创建表格幻灯片"""
        # 解析表头
        header_line = lines[start_idx].strip()
        headers = [cell.strip() for cell in header_line.split('|') if cell.strip()]
        
        # 跳过分隔行
        i = start_idx + 2
        
        # 收集数据行
        rows = []
        while i < len(lines) and '|' in lines[i]:
            row_line = lines[i].strip()
            cells = [cell.strip() for cell in row_line.split('|') if cell.strip()]
            if cells:
                rows.append(cells)
            i += 1
        
        # 创建幻灯片
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "数据表格"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.name = self.default_font
        
        # 添加表格
        if headers:
            num_rows = len(rows) + 1
            num_cols = len(headers)
            
            # 计算表格尺寸
            table_width = Inches(9)
            table_height = Inches(5.5)
            
            table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.2), table_width, table_height).table
            
            # 填充表头
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                cell.text_frame.paragraphs[0].font.name = self.default_font
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            
            # 填充数据
            for row_idx, row_data in enumerate(rows, 1):
                for col_idx, cell_value in enumerate(row_data):
                    if col_idx < num_cols:
                        cell = table.cell(row_idx, col_idx)
                        cell.text = str(cell_value)
                        cell.text_frame.paragraphs[0].font.size = Pt(11)
                        cell.text_frame.paragraphs[0].font.name = self.default_font
        
        return i
    
    def _create_code_slide(self, lines: list, start_idx: int) -> int:
        """创建代码幻灯片"""
        language = lines[start_idx].strip()[3:].strip()
        code_lines = []
        i = start_idx + 1
        
        while i < len(lines) and not lines[i].strip().startswith('```'):
            code_lines.append(lines[i])
            i += 1
        
        code_content = '\n'.join(code_lines)
        
        # 创建幻灯片
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = f"代码示例 {f'({language})' if language else ''}"
        title_frame.paragraphs[0].font.size = Pt(20)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.name = self.default_font
        
        # 添加代码框
        code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(6))
        code_frame = code_box.text_frame
        code_frame.word_wrap = False
        
        # 添加代码内容
        p = code_frame.paragraphs[0]
        p.text = code_content
        p.font.size = self.code_font_size
        p.font.name = 'Consolas'
        
        return i + 1
    
    def _create_image_slide(self, line: str, lines: list, start_idx: int) -> int:
        """创建图片幻灯片"""
        # 解析图片语法 ![alt](path)
        match = re.match(r'!\[([^\]]*)\]\(([^)]+)\)', line.strip())
        if match:
            alt_text = match.group(1)
            image_path = match.group(2)
            
            # 创建幻灯片
            blank_layout = self.prs.slide_layouts[6]
            slide = self.prs.slides.add_slide(blank_layout)
            
            # 添加图片
            if os.path.exists(image_path):
                try:
                    slide.shapes.add_picture(image_path, Inches(2), Inches(1.5), width=Inches(6))
                except:
                    pass
            
            # 添加标题
            if alt_text:
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
                title_frame = title_box.text_frame
                title_frame.text = alt_text
                title_frame.paragraphs[0].font.size = Pt(20)
                title_frame.paragraphs[0].font.bold = True
                title_frame.paragraphs[0].font.name = self.default_font
                title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return start_idx + 1
    
    def _create_paragraph_slide(self, first_line: str, lines: list, start_idx: int) -> int:
        """创建段落幻灯片"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        
        # 收集段落内容
        paragraph_lines = [first_line]
        i = start_idx
        
        while i < len(lines):
            line = lines[i].strip()
            if not line or line.startswith('#') or line.startswith('-') or line.startswith('*'):
                break
            paragraph_lines.append(line)
            i += 1
        
        content = ' '.join(paragraph_lines)
        
        # 添加内容
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.text = content
        content_frame.paragraphs[0].font.size = self.content_font_size
        content_frame.paragraphs[0].font.name = self.default_font
        
        return i


class PPTToMarkdownConverter:
    """PPT 到 Markdown 转换器"""
    
    def __init__(self):
        pass
    
    def convert(self, pptx_path: str, output_path: str = None) -> str:
        """
        将 PPT 文件转换为 Markdown
        
        Args:
            pptx_path: PPT 文件路径
            output_path: 输出 Markdown 文件路径（可选）
        
        Returns:
            Markdown 文本内容
        """
        prs = Presentation(pptx_path)
        md_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            md_content.append(f"## 幻灯片 {slide_num}\n")
            
            # 提取文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    md_content.append(shape.text.strip())
                    md_content.append('')
            
            md_content.append('---\n')
        
        md_text = '\n'.join(md_content)
        
        # 保存到文件
        if output_path:
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(md_text)
        
        return md_text


class PPTImageHandler:
    """PPT 图片处理器 - 支持图文混排"""
    
    def __init__(self, prs=None):
        self.prs = prs
        self.default_font = '微软雅黑'
        
    def set_presentation(self, prs):
        """设置演示文稿对象"""
        self.prs = prs
    
    def insert_image_to_slide(self, slide_index: int, image_path: str, 
                              left: float = 1, top: float = 1, 
                              width: float = 4, height: float = None) -> bool:
        """
        在指定幻灯片插入图片
        
        Args:
            slide_index: 幻灯片索引（从0开始）
            image_path: 图片路径
            left: 左边距（英寸）
            top: 上边距（英寸）
            width: 图片宽度（英寸）
            height: 图片高度（英寸，可选）
        
        Returns:
            是否成功
        """
        try:
            if not os.path.exists(image_path):
                print(f"图片不存在: {image_path}")
                return False
            
            if slide_index < 0 or slide_index >= len(self.prs.slides):
                print(f"幻灯片索引超出范围: {slide_index}")
                return False
            
            slide = self.prs.slides[slide_index]
            
            # 插入图片
            if height:
                slide.shapes.add_picture(image_path, Inches(left), Inches(top), 
                                        Inches(width), Inches(height))
            else:
                slide.shapes.add_picture(image_path, Inches(left), Inches(top), 
                                        Inches(width))
            
            return True
            
        except Exception as e:
            print(f"插入图片失败: {e}")
            return False
    
    def create_text_image_layout(self, slide_index: int, text: str, image_path: str,
                                  layout: str = 'left', image_width: float = 4) -> bool:
        """
        在指定幻灯片创建图文混排布局
        
        Args:
            slide_index: 幻灯片索引（从0开始）
            text: 文本内容
            image_path: 图片路径
            layout: 布局方式 (left/right/top/bottom)
            image_width: 图片宽度（英寸）
        
        Returns:
            是否成功
        """
        try:
            if not os.path.exists(image_path):
                print(f"图片不存在: {image_path}")
                return False
            
            if slide_index < 0 or slide_index >= len(self.prs.slides):
                print(f"幻灯片索引超出范围: {slide_index}")
                return False
            
            slide = self.prs.slides[slide_index]
            
            # 幻灯片尺寸
            slide_width = 10  # 英寸
            slide_height = 7.5  # 英寸
            margin = 0.5  # 边距
            
            if layout == 'left':
                # 图片左，文字右
                img_left = margin
                img_top = margin + 1
                text_left = margin + image_width + 0.5
                text_top = margin + 1
                text_width = slide_width - text_left - margin
                text_height = slide_height - text_top - margin
                
            elif layout == 'right':
                # 图片右，文字左
                text_left = margin
                text_top = margin + 1
                text_width = slide_width - margin - image_width - margin - 0.5
                text_height = slide_height - text_top - margin
                img_left = slide_width - margin - image_width
                img_top = margin + 1
                
            elif layout == 'top':
                # 图片上，文字下
                img_left = (slide_width - image_width) / 2
                img_top = margin + 1
                text_left = margin
                text_top = img_top + 3  # 假设图片高度约3英寸
                text_width = slide_width - 2 * margin
                text_height = slide_height - text_top - margin
                
            elif layout == 'bottom':
                # 文字上，图片下
                text_left = margin
                text_top = margin + 1
                text_width = slide_width - 2 * margin
                text_height = 2.5
                img_left = (slide_width - image_width) / 2
                img_top = text_top + text_height + 0.5
                
            else:
                print(f"不支持的布局方式: {layout}")
                return False
            
            # 插入图片
            try:
                slide.shapes.add_picture(image_path, Inches(img_left), Inches(img_top), 
                                        Inches(image_width))
            except Exception as e:
                print(f"插入图片失败: {e}")
                return False
            
            # 添加文本框
            text_box = slide.shapes.add_textbox(Inches(text_left), Inches(text_top), 
                                               Inches(text_width), Inches(text_height))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.text = text
            
            # 设置字体
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.name = self.default_font
            
            return True
            
        except Exception as e:
            print(f"创建图文混排失败: {e}")
            return False


# 便捷函数
def md_to_ppt(md_file: str, output_path: str = None, title: str = None) -> str:
    """Markdown 文件转 PPT"""
    converter = MarkdownToPPTConverter()
    return converter.convert_file(md_file, output_path, title)


def ppt_to_md(pptx_file: str, output_path: str = None) -> str:
    """PPT 文件转 Markdown"""
    converter = PPTToMarkdownConverter()
    return converter.convert(pptx_file, output_path)


def insert_image_to_ppt(pptx_file: str, slide_index: int, image_path: str,
                        output_path: str = None, **kwargs) -> str:
    """
    向 PPT 指定幻灯片插入图片
    
    Args:
        pptx_file: PPT 文件路径
        slide_index: 幻灯片索引（从0开始）
        image_path: 图片路径
        output_path: 输出路径（可选）
        **kwargs: 其他参数（left, top, width, height）
    
    Returns:
        输出文件路径
    """
    from pptx import Presentation
    
    prs = Presentation(pptx_file)
    handler = PPTImageHandler(prs)
    
    handler.insert_image_to_slide(
        slide_index,
        image_path,
        left=kwargs.get('left', 1),
        top=kwargs.get('top', 1),
        width=kwargs.get('width', 4),
        height=kwargs.get('height', None)
    )
    
    if output_path is None:
        output_path = pptx_file
    
    prs.save(output_path)
    return output_path


def create_ppt_text_image_layout(pptx_file: str, slide_index: int, text: str,
                                  image_path: str, layout: str = 'left',
                                  image_width: float = 4, output_path: str = None) -> str:
    """
    在 PPT 指定幻灯片创建图文混排布局
    
    Args:
        pptx_file: PPT 文件路径
        slide_index: 幻灯片索引（从0开始）
        text: 文本内容
        image_path: 图片路径
        layout: 布局方式 (left/right/top/bottom)
        image_width: 图片宽度（英寸）
        output_path: 输出路径（可选）
    
    Returns:
        输出文件路径
    """
    from pptx import Presentation
    
    prs = Presentation(pptx_file)
    handler = PPTImageHandler(prs)
    
    handler.create_text_image_layout(slide_index, text, image_path, layout, image_width)
    
    if output_path is None:
        output_path = pptx_file
    
    prs.save(output_path)
    return output_path


if __name__ == '__main__':
    # 测试代码
    test_md = """# 项目汇报

## 项目概述

本项目旨在开发一款智能文档处理工具。

## 核心功能

- 文档自动转换
- 智能格式识别
- 批量处理能力

## 技术架构

| 模块 | 技术栈 | 状态 |
|------|--------|------|
| 前端 | React | 完成 |
| 后端 | Python | 完成 |
| 数据库 | PostgreSQL | 完成 |

## 代码示例

```python
def hello():
    print("Hello World")
```

## 总结

项目进展顺利，预计下月上线。
"""
    
    converter = MarkdownToPPTConverter()
    output = '/tmp/test_ppt.pptx'
    converter.convert(test_md, output, '项目汇报')
    print(f'转换完成: {output}')
