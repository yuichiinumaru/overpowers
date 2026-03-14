#!/usr/bin/env python3
"""
生成教案 PPT - 支持多种课文类型
用法: python generate_ppt.py <课文名> <年级> <课文类型> [备课资料文件路径]
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_ppt(lesson_name, grade, lesson_type, content_file=None):
    """基于备课资料生成教案 PPT"""
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 根据课文类型调整副标题
    type_names = {
        '古诗': '古诗教学',
        '古诗词': '古诗教学',
        '现代文': '现代文教学',
        '散文': '散文教学',
        '寓言': '寓言教学',
        '童话': '童话教学',
        '说明文': '说明文教学'
    }
    type_display = type_names.get(lesson_type, '语文教学')
    
    # 第1页：标题页
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加标题背景
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
                                    prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(70, 130, 180)
    shape.line.fill.background()
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"《{lesson_name}》"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    p2 = tf.add_paragraph()
    p2.text = f"{grade} {type_display}教案"
    p2.font.size = Pt(32)
    p2.font.color.rgb = RGBColor(255, 223, 186)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)
    
    # 第2页：教学目标
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "教学目标"
    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    tf.text = "知识目标："
    p = tf.add_paragraph()
    p.text = "• 认识本课生字，正确、流利地朗读课文"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 理解课文内容，掌握重点词语"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "能力目标："
    p = tf.add_paragraph()
    p.text = "• 培养阅读理解能力和语言表达能力"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "情感目标："
    p = tf.add_paragraph()
    p.text = "• 体会课文情感，激发学习兴趣"
    p.level = 1
    
    # 第3页：教学重难点
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "教学重难点"
    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    tf.text = "教学重点："
    p = tf.add_paragraph()
    p.text = "• 理解课文主要内容"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 掌握重点字词，正确流利朗读"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "教学难点："
    p = tf.add_paragraph()
    p.text = "• 体会课文表达的思想感情"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 学习作者的表达方法"
    p.level = 1
    
    # 第4页：教学过程-导入
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "教学过程 - 导入新课"
    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    tf.text = "一、创设情境，导入新课（5分钟）"
    p = tf.add_paragraph()
    p.text = "• 谈话导入，激发兴趣"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 揭示课题，齐读课题"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 简介作者/背景（如需要）"
    p.level = 1
    
    # 第5页：教学过程-初读
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "教学过程 - 初读感知"
    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    tf.text = "二、初读课文，整体感知（10分钟）"
    p = tf.add_paragraph()
    p.text = "• 自由朗读课文，标出生字词"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 检查生字词认读情况"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 指名分段朗读课文"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 思考：课文主要讲了什么？"
    p.level = 1
    
    # 第6页：教学过程-精读
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "教学过程 - 精读理解"
    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    tf.text = "三、精读课文，深入理解（15分钟）"
    p = tf.add_paragraph()
    p.text = "• 品读重点段落，理解关键词句"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 讨论交流，体会情感"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 指导有感情地朗读"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 学习写作方法（如适用）"
    p.level = 1
    
    # 第7页：教学过程-巩固
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "教学过程 - 巩固拓展"
    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    tf.text = "四、总结提升，拓展延伸（10分钟）"
    p = tf.add_paragraph()
    p.text = "• 回顾课文主要内容"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 总结学习收获"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 拓展阅读/相关练习"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 布置课后作业"
    p.level = 1
    
    # 第8页：板书设计
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "板书设计"
    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    tf.text = f"《{lesson_name}》"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "【主要内容】"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "【重点字词】"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "【中心思想】"
    
    # 第9页：作业布置
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "作业布置"
    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    tf.text = "1. 熟读课文，背诵重点段落（如需要）"
    p = tf.add_paragraph()
    p.text = "2. 抄写生字词，完成课后习题"
    p = tf.add_paragraph()
    p.text = "3. 预习下一课内容"
    p = tf.add_paragraph()
    if lesson_type in ['古诗', '古诗词']:
        p = tf.add_paragraph()
        p.text = "4. 搜集其他相关古诗，读一读"
    else:
        p = tf.add_paragraph()
        p.text = "4. 观察生活中的相关现象，写几句话"
    
    # 保存
    output_file = f"教案_{lesson_name}.pptx"
    prs.save(output_file)
    print(f"PPT 已生成: {output_file}")
    return output_file

if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("用法: python generate_ppt.py <课文名> <年级> <课文类型> [备课资料文件路径]")
        print("课文类型: 古诗/现代文/寓言/童话/说明文")
        sys.exit(1)
    
    lesson_name = sys.argv[1]
    grade = sys.argv[2]
    lesson_type = sys.argv[3]
    content_file = sys.argv[4] if len(sys.argv) > 4 else None
    
    create_ppt(lesson_name, grade, lesson_type, content_file)
