#!/usr/bin/env python3
"""
Auto-PPT Generator — takes structured JSON and produces a clean PPTX (+ PDF).

Usage:
    python3 make_ppt.py slides.json [--out /path/to/output.pptx] [--pdf]

JSON schema:
{
  "title": "Presentation Title",
  "author": "",
  "slides": [
    { "type": "cover",      "title": "...", "subtitle": "..." },
    { "type": "content",    "title": "...", "bullets": ["...", "..."] },
    { "type": "two_column", "title": "...", "left_title": "A", "left": ["..."], "right_title": "B", "right": ["..."] },
    { "type": "flow",       "title": "...", "steps": ["Step 1", "Step 2", "Step 3"] },
    { "type": "comparison", "title": "...", "headers": ["A","B"], "rows": [["a1","b1"],["a2","b2"]] },
    { "type": "summary",    "title": "...", "bullets": ["..."] },
    { "type": "end",        "title": "Thank You", "subtitle": "..." }
  ]
}
"""

import json
import sys
import os
import argparse
import subprocess

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

FONT_TITLE = "Microsoft YaHei"
FONT_BODY = "Microsoft YaHei"
FONT_FALLBACK = "PingFang SC"

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xEE, 0xEE, 0xEE)
ACCENT = RGBColor(0x00, 0x78, 0xD4)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.8)
MARGIN_T = Inches(0.6)
CONTENT_W = Inches(11.7)


def _set_font(run, name=FONT_BODY, size=Pt(18), bold=False, color=BLACK):
    run.font.name = name
    run.font.east_asian_font = FONT_FALLBACK
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color


def _set_slide_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_text(slide, text, left=MARGIN_L, top=MARGIN_T, width=CONTENT_W,
                    height=Inches(0.9), size=Pt(40), bold=True, color=BLACK,
                    align=PP_ALIGN.LEFT):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, name=FONT_TITLE, size=size, bold=bold, color=color)
    return txbox


def _add_underline(slide, left, top, width, color=ACCENT, height=Pt(3)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_bullets(slide, items, left=MARGIN_L, top=Inches(2.0),
                 width=CONTENT_W, height=Inches(4.5),
                 font_size=Pt(22), color=DARK_GRAY):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"  {item}"
        _set_font(run, size=font_size, color=color)
        p.level = 0
    return txbox


def _add_page_number(slide, num, total):
    txbox = slide.shapes.add_textbox(
        Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.4)
    )
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{num} / {total}"
    _set_font(run, size=Pt(11), color=MID_GRAY)


def build_cover(prs, data, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide)

    _add_title_text(slide, data.get("title", ""),
                    left=MARGIN_L, top=Inches(2.2), width=CONTENT_W,
                    height=Inches(1.2), size=Pt(44), bold=True, color=BLACK,
                    align=PP_ALIGN.LEFT)

    _add_underline(slide, MARGIN_L, Inches(3.5), Inches(2.0))

    subtitle = data.get("subtitle", "")
    if subtitle:
        txbox = slide.shapes.add_textbox(MARGIN_L, Inches(3.8), CONTENT_W, Inches(0.7))
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = subtitle
        _set_font(run, size=Pt(20), color=MID_GRAY)


def build_content(prs, data, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    _add_title_text(slide, data.get("title", ""))
    _add_underline(slide, MARGIN_L, Inches(1.5), Inches(1.5))
    _add_bullets(slide, data.get("bullets", []))
    _add_page_number(slide, page_num, total)


def build_two_column(prs, data, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    _add_title_text(slide, data.get("title", ""))
    _add_underline(slide, MARGIN_L, Inches(1.5), Inches(1.5))

    col_w = Inches(5.3)
    gap = Inches(0.6)

    lt = data.get("left_title", "")
    if lt:
        txbox = slide.shapes.add_textbox(MARGIN_L, Inches(1.8), col_w, Inches(0.5))
        tf = txbox.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = lt
        _set_font(run, size=Pt(22), bold=True, color=ACCENT)

    _add_bullets(slide, data.get("left", []),
                 left=MARGIN_L, top=Inches(2.3), width=col_w, height=Inches(4.0),
                 font_size=Pt(18))

    right_left = MARGIN_L + col_w + gap
    rt = data.get("right_title", "")
    if rt:
        txbox = slide.shapes.add_textbox(right_left, Inches(1.8), col_w, Inches(0.5))
        tf = txbox.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = rt
        _set_font(run, size=Pt(22), bold=True, color=ACCENT)

    _add_bullets(slide, data.get("right", []),
                 left=right_left, top=Inches(2.3), width=col_w, height=Inches(4.0),
                 font_size=Pt(18))

    _add_page_number(slide, page_num, total)


def build_flow(prs, data, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    _add_title_text(slide, data.get("title", ""))
    _add_underline(slide, MARGIN_L, Inches(1.5), Inches(1.5))

    steps = data.get("steps", [])
    n = len(steps)
    if n == 0:
        _add_page_number(slide, page_num, total)
        return

    usable_w = CONTENT_W
    box_w = min(Inches(2.2), usable_w / max(n, 1))
    gap = Inches(0.3) if n > 1 else Inches(0)
    total_w = n * box_w + (n - 1) * gap
    start_x = MARGIN_L + (usable_w - total_w) / 2
    box_y = Inches(3.2)
    box_h = Inches(1.6)
    arrow_y = box_y + box_h / 2 - Pt(10)

    for i, step in enumerate(steps):
        x = start_x + i * (box_w + gap)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(box_y), int(box_w), int(box_h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT
        shape.line.fill.background()
        shape.shadow.inherit = False

        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = step
        _set_font(run, size=Pt(16), bold=True, color=WHITE)
        tf.paragraphs[0].space_before = Pt(4)

        num_p = tf.add_paragraph()
        num_p.alignment = PP_ALIGN.CENTER
        num_run = num_p.add_run()
        num_run.text = f"Step {i + 1}"
        _set_font(num_run, size=Pt(11), color=RGBColor(0xCC, 0xE5, 0xFF))

        if i < n - 1:
            ax = x + box_w + Pt(4)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, int(ax), int(arrow_y),
                int(gap - Pt(8)), Pt(20)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MID_GRAY
            arrow.line.fill.background()

    _add_page_number(slide, page_num, total)


def build_comparison(prs, data, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    _add_title_text(slide, data.get("title", ""))
    _add_underline(slide, MARGIN_L, Inches(1.5), Inches(1.5))

    headers = data.get("headers", [])
    rows = data.get("rows", [])
    n_cols = len(headers) if headers else (len(rows[0]) if rows else 0)
    n_rows = len(rows) + (1 if headers else 0)

    if n_cols == 0 or n_rows == 0:
        _add_page_number(slide, page_num, total)
        return

    tbl_left = MARGIN_L
    tbl_top = Inches(2.0)
    tbl_w = CONTENT_W
    tbl_h = Inches(0.55) * n_rows

    table_shape = slide.shapes.add_table(n_rows, n_cols, int(tbl_left), int(tbl_top), int(tbl_w), int(tbl_h))
    table = table_shape.table

    col_w = int(tbl_w / n_cols)
    for ci in range(n_cols):
        table.columns[ci].width = col_w

    if headers:
        for ci, h in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = h
            _set_font(run, size=Pt(18), bold=True, color=WHITE)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT

    for ri, row in enumerate(rows):
        tr = ri + (1 if headers else 0)
        for ci, val in enumerate(row):
            cell = table.cell(tr, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(val)
            _set_font(run, size=Pt(16), color=DARK_GRAY)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 == 0 else LIGHT_GRAY

    _add_page_number(slide, page_num, total)


def build_summary(prs, data, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    _add_title_text(slide, data.get("title", "Summary"),
                    size=Pt(40), bold=True, color=BLACK)
    _add_underline(slide, MARGIN_L, Inches(1.5), Inches(1.5), color=ACCENT)
    _add_bullets(slide, data.get("bullets", []), font_size=Pt(24), color=DARK_GRAY)
    _add_page_number(slide, page_num, total)


def build_end(prs, data, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    _add_title_text(slide, data.get("title", "Thank You"),
                    left=MARGIN_L, top=Inches(2.6), width=CONTENT_W,
                    height=Inches(1.0), size=Pt(44), bold=True, color=BLACK,
                    align=PP_ALIGN.CENTER)

    sub = data.get("subtitle", "")
    if sub:
        txbox = slide.shapes.add_textbox(MARGIN_L, Inches(3.8), CONTENT_W, Inches(0.6))
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = sub
        _set_font(run, size=Pt(20), color=MID_GRAY)

    _add_page_number(slide, total, total)


BUILDERS = {
    "cover": lambda prs, d, pn, t: build_cover(prs, d, t),
    "content": build_content,
    "two_column": build_two_column,
    "flow": build_flow,
    "comparison": build_comparison,
    "summary": build_summary,
    "end": lambda prs, d, pn, t: build_end(prs, d, t),
}


def generate(json_path, out_path, to_pdf=False):
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    slides = data.get("slides", [])
    total = len(slides)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for i, s in enumerate(slides):
        stype = s.get("type", "content")
        builder = BUILDERS.get(stype, build_content)
        builder(prs, s, i + 1, total)

    prs.save(out_path)
    print(f"PPTX saved: {out_path}")

    if to_pdf:
        pdf_path = out_path.rsplit(".", 1)[0] + ".pdf"
        converted = False

        for cmd in [
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir",
             os.path.dirname(os.path.abspath(out_path)), out_path],
            ["/Applications/LibreOffice.app/Contents/MacOS/soffice",
             "--headless", "--convert-to", "pdf", "--outdir",
             os.path.dirname(os.path.abspath(out_path)), out_path],
        ]:
            try:
                subprocess.run(cmd, check=True, capture_output=True, timeout=60)
                converted = True
                break
            except (FileNotFoundError, subprocess.CalledProcessError):
                continue

        if not converted:
            try:
                applescript = f'''
                tell application "Keynote"
                    set theDoc to open POSIX file "{os.path.abspath(out_path)}"
                    delay 3
                    export theDoc to POSIX file "{os.path.abspath(pdf_path)}" as PDF
                    close theDoc saving no
                end tell
                '''
                subprocess.run(["osascript", "-e", applescript],
                               check=True, capture_output=True, timeout=30)
                converted = True
            except (FileNotFoundError, subprocess.CalledProcessError):
                pass

        if converted and os.path.exists(pdf_path):
            print(f"PDF  saved: {pdf_path}")
        else:
            print(f"PDF conversion skipped (install LibreOffice or use Keynote). PPTX is ready.")


def main():
    parser = argparse.ArgumentParser(description="Generate PPT from structured JSON")
    parser.add_argument("json_file", help="Path to the JSON slide definition file")
    parser.add_argument("--out", "-o", default=None, help="Output PPTX path (default: ~/Desktop/<title>.pptx)")
    parser.add_argument("--pdf", action="store_true", help="Also convert to PDF")
    args = parser.parse_args()

    if not os.path.exists(args.json_file):
        print(f"Error: {args.json_file} not found", file=sys.stderr)
        sys.exit(1)

    with open(args.json_file, "r", encoding="utf-8") as f:
        data = json.load(f)

    out = args.out
    if not out:
        title = data.get("title", "presentation").replace(" ", "_").replace("/", "_")
        out = os.path.expanduser(f"~/Desktop/{title}.pptx")

    generate(args.json_file, out, to_pdf=args.pdf)


if __name__ == "__main__":
    main()
