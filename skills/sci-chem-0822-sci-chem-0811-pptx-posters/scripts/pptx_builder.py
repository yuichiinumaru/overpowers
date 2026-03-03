from pptx import Presentation
from pptx.util import Inches, Pt
import os
import sys

def create_pptx_poster(output_file="poster.pptx", width_inches=48, height_inches=36):
    """Create a blank PPTX poster with specified dimensions"""
    prs = Presentation()
    
    # Set slide dimensions
    prs.slide_width = Inches(width_inches)
    prs.slide_height = Inches(height_inches)
    
    # Add a blank slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    prs.save(output_file)
    print(f"Created blank poster: {output_file} ({width_inches}x{height_inches} inches)")
    return prs, slide

def add_image_to_poster(pptx_file, image_path, left_inch, top_inch, width_inch=None):
    """Add an image to an existing PPTX poster"""
    if not os.path.exists(pptx_file):
        print(f"Error: {pptx_file} not found.")
        return
        
    prs = Presentation(pptx_file)
    slide = prs.slides[0]
    
    if width_inch:
        slide.shapes.add_picture(image_path, Inches(left_inch), Inches(top_inch), width=Inches(width_inch))
    else:
        slide.shapes.add_picture(image_path, Inches(left_inch), Inches(top_inch))
        
    prs.save(pptx_file)
    print(f"Added {image_path} to {pptx_file}")

if __name__ == "__main__":
    if len(sys.argv) > 1 and sys.argv[1] == "create":
        create_pptx_poster()
    else:
        print("Usage:")
        print("  python pptx_builder.py create")
        print("  (This is a helper module for programmatic PPTX generation)")
        
