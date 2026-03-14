#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
OpenClaw 知识教学技能 - 文档生成脚本

功能:
- 从知识库生成 PPT 演示文稿
- 从知识库生成 Word 文档
- 支持自定义模板和样式
"""

import argparse
import os
import re
import sys
from datetime import datetime
from pathlib import Path

# 添加父目录到路径以导入模块
sys.path.insert(0, str(Path(__file__).parent.parent.parent / 'docx' / 'scripts'))

try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.style import WD_STYLE_TYPE
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    print("警告: python-docx 未安装，Word 文档生成功能不可用")

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    from pptx.dml.color import RGBColor as PptxRGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("警告: python-pptx 未安装，PPT 生成功能不可用")


class KnowledgeBaseParser:
    """知识库解析器"""
    
    def __init__(self, knowledge_file: str):
        self.knowledge_file = knowledge_file
        self.content = self._load_content()
    
    def _load_content(self) -> str:
        """加载知识库内容"""
        with open(self.knowledge_file, 'r', encoding='utf-8') as f:
            return f.read()
    
    def get_sections(self) -> list:
        """获取所有章节"""
        sections = []
        pattern = r'^## (\d+)\. (.+)$'
        
        for match in re.finditer(pattern, self.content, re.MULTILINE):
            sections.append({
                'number': match.group(1),
                'title': match.group(2),
                'start': match.end()
            })
        
        # 添加内容到每个章节
        for i, section in enumerate(sections):
            end = sections[i + 1]['start'] if i + 1 < len(sections) else len(self.content)
            section['content'] = self.content[section['start']:end].strip()
        
        return sections
    
    def get_subsections(self, section_content: str) -> list:
        """获取子章节"""
        subsections = []
        pattern = r'^### (\d+\.\d+) (.+)$'
        
        for match in re.finditer(pattern, section_content, re.MULTILINE):
            subsections.append({
                'number': match.group(1),
                'title': match.group(2)
            })
        
        return subsections
    
    def get_topic_content(self, topic: str) -> dict:
        """获取特定主题的内容"""
        sections = self.get_sections()
        
        for section in sections:
            if topic.lower() in section['title'].lower():
                return section
        
        # 如果没找到精确匹配，返回相关内容
        for section in sections:
            if topic.lower() in section['content'].lower():
                return section
        
        return None


class PPTGenerator:
    """PPT 生成器"""
    
    def __init__(self, template: str = 'modern'):
        self.template = template
        self.prs = Presentation()
        self._setup_slide_size()
    
    def _setup_slide_size(self):
        """设置幻灯片尺寸"""
        self.prs.slide_width = PptxInches(13.333)
        self.prs.slide_height = PptxInches(7.5)
    
    def _get_colors(self) -> dict:
        """获取模板颜色"""
        templates = {
            'modern': {
                'primary': PptxRGBColor(11, 18, 32),      # 深蓝黑
                'secondary': PptxRGBColor(42, 43, 43),    # 深灰
                'accent': PptxRGBColor(154, 166, 178),    # 银灰
                'background': PptxRGBColor(255, 255, 255) # 白色
            },
            'professional': {
                'primary': PptxRGBColor(0, 51, 102),      # 深蓝
                'secondary': PptxRGBColor(51, 102, 153),  # 中蓝
                'accent': PptxRGBColor(255, 153, 0),      # 橙色
                'background': PptxRGBColor(255, 255, 255)
            },
            'creative': {
                'primary': PptxRGBColor(102, 51, 153),    # 紫色
                'secondary': PptxRGBColor(51, 153, 102),  # 绿色
                'accent': PptxRGBColor(255, 102, 102),    # 红色
                'background': PptxRGBColor(250, 250, 250)
            }
        }
        return templates.get(self.template, templates['modern'])
    
    def add_title_slide(self, title: str, subtitle: str = ''):
        """添加标题幻灯片"""
        slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(slide_layout)
        
        colors = self._get_colors()
        
        # 添加标题
        title_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(2.5),
            PptxInches(12.333), PptxInches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = PptxPt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = colors['primary']
        title_para.alignment = PP_ALIGN.CENTER
        
        # 添加副标题
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                PptxInches(0.5), PptxInches(4.2),
                PptxInches(12.333), PptxInches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = subtitle
            subtitle_para.font.size = PptxPt(24)
            subtitle_para.font.color.rgb = colors['secondary']
            subtitle_para.alignment = PP_ALIGN.CENTER
    
    def add_section_slide(self, section_number: str, title: str):
        """添加章节标题幻灯片"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        colors = self._get_colors()
        
        # 章节编号
        num_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(2),
            PptxInches(12.333), PptxInches(1)
        )
        num_frame = num_box.text_frame
        num_para = num_frame.paragraphs[0]
        num_para.text = f"第 {section_number} 章"
        num_para.font.size = PptxPt(28)
        num_para.font.color.rgb = colors['accent']
        num_para.alignment = PP_ALIGN.CENTER
        
        # 章节标题
        title_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(3),
            PptxInches(12.333), PptxInches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = PptxPt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = colors['primary']
        title_para.alignment = PP_ALIGN.CENTER
    
    def add_content_slide(self, title: str, content: str):
        """添加内容幻灯片"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        colors = self._get_colors()
        
        # 标题
        title_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.3),
            PptxInches(12.333), PptxInches(0.8)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = PptxPt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = colors['primary']
        
        # 内容
        content_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(1.3),
            PptxInches(12.333), PptxInches(5.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        # 处理内容
        lines = content.split('\n')
        for i, line in enumerate(lines[:15]):  # 限制行数
            if i == 0:
                para = content_frame.paragraphs[0]
            else:
                para = content_frame.add_paragraph()
            
            # 清理 markdown 格式
            clean_line = re.sub(r'[#*`]', '', line).strip()
            if clean_line:
                para.text = clean_line
                para.font.size = PptxPt(18)
                para.font.color.rgb = colors['secondary']
                para.space_after = PptxPt(8)
    
    def add_toc_slide(self, sections: list):
        """添加目录幻灯片"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        colors = self._get_colors()
        
        # 标题
        title_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.5),
            PptxInches(12.333), PptxInches(1)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "目录"
        title_para.font.size = PptxPt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = colors['primary']
        title_para.alignment = PP_ALIGN.CENTER
        
        # 目录项
        content_box = slide.shapes.add_textbox(
            PptxInches(1), PptxInches(1.8),
            PptxInches(11.333), PptxInches(5)
        )
        content_frame = content_box.text_frame
        
        for i, section in enumerate(sections):
            if i == 0:
                para = content_frame.paragraphs[0]
            else:
                para = content_frame.add_paragraph()
            
            para.text = f"{section['number']}. {section['title']}"
            para.font.size = PptxPt(22)
            para.font.color.rgb = colors['secondary']
            para.space_after = PptxPt(12)
    
    def save(self, output_path: str):
        """保存 PPT 文件"""
        self.prs.save(output_path)


class WordGenerator:
    """Word 文档生成器"""
    
    def __init__(self, style: str = 'tutorial'):
        self.style = style
        self.doc = Document()
        self._setup_styles()
    
    def _setup_styles(self):
        """设置文档样式"""
        # 设置默认字体
        style = self.doc.styles['Normal']
        font = style.font
        font.name = 'SimHei'
        font.size = Pt(12)
    
    def add_cover(self, title: str, subtitle: str = ''):
        """添加封面"""
        # 标题
        title_para = self.doc.add_paragraph()
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title_run = title_para.add_run(title)
        title_run.font.size = Pt(28)
        title_run.font.bold = True
        
        # 副标题
        if subtitle:
            subtitle_para = self.doc.add_paragraph()
            subtitle_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            subtitle_run = subtitle_para.add_run(subtitle)
            subtitle_run.font.size = Pt(16)
        
        # 日期
        date_para = self.doc.add_paragraph()
        date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        date_run = date_para.add_run(datetime.now().strftime('%Y-%m-%d'))
        date_run.font.size = Pt(12)
        
        self.doc.add_page_break()
    
    def add_toc(self):
        """添加目录占位符"""
        toc_para = self.doc.add_paragraph()
        toc_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        toc_run = toc_para.add_run('目录')
        toc_run.font.size = Pt(20)
        toc_run.font.bold = True
        
        self.doc.add_paragraph()
        self.doc.add_page_break()
    
    def add_heading(self, text: str, level: int = 1):
        """添加标题"""
        heading = self.doc.add_heading(text, level=level)
        return heading
    
    def add_paragraph(self, text: str):
        """添加段落"""
        # 清理 markdown 格式
        clean_text = re.sub(r'[#*`]', '', text).strip()
        if clean_text:
            para = self.doc.add_paragraph(clean_text)
            return para
        return None
    
    def add_code_block(self, code: str, language: str = ''):
        """添加代码块"""
        para = self.doc.add_paragraph()
        para.style = 'No Spacing'
        run = para.add_run(code)
        run.font.name = 'Courier New'
        run.font.size = Pt(10)
    
    def add_table(self, headers: list, rows: list):
        """添加表格"""
        table = self.doc.add_table(rows=len(rows) + 1, cols=len(headers))
        table.style = 'Table Grid'
        
        # 添加表头
        for i, header in enumerate(headers):
            cell = table.rows[0].cells[i]
            cell.text = header
        
        # 添加数据行
        for i, row in enumerate(rows):
            for j, cell_text in enumerate(row):
                table.rows[i + 1].cells[j].text = cell_text
    
    def save(self, output_path: str):
        """保存文档"""
        self.doc.save(output_path)


def generate_ppt(knowledge_file: str, output_path: str, topic: str = None, template: str = 'modern'):
    """生成 PPT 文档"""
    if not HAS_PPTX:
        print("错误: python-pptx 未安装，无法生成 PPT")
        return False
    
    parser = KnowledgeBaseParser(knowledge_file)
    generator = PPTGenerator(template=template)
    
    # 添加标题幻灯片
    generator.add_title_slide(
        title='OpenClaw 知识教程',
        subtitle='OpenClaw 平台使用指南与技能开发教程'
    )
    
    # 获取章节
    sections = parser.get_sections()
    
    # 如果指定了主题，只生成相关内容
    if topic:
        topic_content = parser.get_topic_content(topic)
        if topic_content:
            sections = [topic_content]
    
    # 添加目录
    generator.add_toc_slide(sections)
    
    # 添加各章节内容
    for section in sections:
        # 章节标题页
        generator.add_section_slide(section['number'], section['title'])
        
        # 章节内容页
        content_lines = section['content'].split('\n')
        current_content = []
        current_title = section['title']
        
        for line in content_lines:
            if line.startswith('### '):
                # 新子章节，保存之前的内容
                if current_content:
                    generator.add_content_slide(
                        current_title,
                        '\n'.join(current_content)
                    )
                current_title = re.sub(r'^### \d+\.\d+ ', '', line)
                current_content = []
            elif line.startswith('#### '):
                # 四级标题作为内容的一部分
                current_content.append(re.sub(r'^#### ', '', line))
            else:
                current_content.append(line)
        
        # 保存最后一部分内容
        if current_content:
            generator.add_content_slide(current_title, '\n'.join(current_content))
    
    # 添加结束页
    generator.add_title_slide(
        title='感谢观看',
        subtitle='OpenClaw - 开放式 AI 技能平台'
    )
    
    # 保存文件
    generator.save(output_path)
    print(f"PPT 已生成: {output_path}")
    return True


def generate_docx(knowledge_file: str, output_path: str, topic: str = None, 
                  include_toc: bool = True, include_cover: bool = True):
    """生成 Word 文档"""
    if not HAS_DOCX:
        print("错误: python-docx 未安装，无法生成 Word 文档")
        return False
    
    parser = KnowledgeBaseParser(knowledge_file)
    generator = WordGenerator()
    
    # 添加封面
    if include_cover:
        generator.add_cover(
            title='OpenClaw 知识教程',
            subtitle='OpenClaw 平台使用指南与技能开发教程'
        )
    
    # 添加目录
    if include_toc:
        generator.add_toc()
    
    # 获取章节
    sections = parser.get_sections()
    
    # 如果指定了主题，只生成相关内容
    if topic:
        topic_content = parser.get_topic_content(topic)
        if topic_content:
            sections = [topic_content]
    
    # 添加各章节内容
    for section in sections:
        # 章节标题
        generator.add_heading(f"{section['number']}. {section['title']}", level=1)
        
        # 处理章节内容
        content_lines = section['content'].split('\n')
        in_code_block = False
        code_content = []
        
        for line in content_lines:
            # 跳过章节标题（已处理）
            if line.startswith('## '):
                continue
            
            # 处理子章节标题
            if line.startswith('### '):
                title = re.sub(r'^### \d+\.\d+ ', '', line)
                generator.add_heading(title, level=2)
                continue
            
            if line.startswith('#### '):
                title = re.sub(r'^#### ', '', line)
                generator.add_heading(title, level=3)
                continue
            
            # 处理代码块
            if line.startswith('```'):
                if in_code_block:
                    # 结束代码块
                    if code_content:
                        generator.add_code_block('\n'.join(code_content))
                    code_content = []
                    in_code_block = False
                else:
                    # 开始代码块
                    in_code_block = True
                continue
            
            if in_code_block:
                code_content.append(line)
                continue
            
            # 处理表格（简化处理）
            if line.startswith('|') and '|' in line[1:]:
                # 跳过表格分隔行
                if re.match(r'^\|[-:\s|]+\|$', line):
                    continue
                # 表格处理可以扩展
                generator.add_paragraph(line)
                continue
            
            # 普通段落
            generator.add_paragraph(line)
    
    # 保存文件
    generator.save(output_path)
    print(f"Word 文档已生成: {output_path}")
    return True


def main():
    parser = argparse.ArgumentParser(description='OpenClaw 知识文档生成工具')
    subparsers = parser.add_subparsers(dest='command', help='可用命令')
    
    # PPT 生成命令
    ppt_parser = subparsers.add_parser('ppt', help='生成 PPT 演示文稿')
    ppt_parser.add_argument('--knowledge', '-k', 
                           default=os.path.join(os.path.dirname(__file__), '..', 'KNOWLEDGE_BASE.md'),
                           help='知识库文件路径')
    ppt_parser.add_argument('--output', '-o', required=True, help='输出文件路径')
    ppt_parser.add_argument('--topic', '-t', help='特定主题（可选）')
    ppt_parser.add_argument('--template', choices=['modern', 'professional', 'creative'],
                           default='modern', help='PPT 模板风格')
    
    # Word 文档生成命令
    docx_parser = subparsers.add_parser('docx', help='生成 Word 文档')
    docx_parser.add_argument('--knowledge', '-k',
                            default=os.path.join(os.path.dirname(__file__), '..', 'KNOWLEDGE_BASE.md'),
                            help='知识库文件路径')
    docx_parser.add_argument('--output', '-o', required=True, help='输出文件路径')
    docx_parser.add_argument('--topic', '-t', help='特定主题（可选）')
    docx_parser.add_argument('--include-toc', action='store_true', help='包含目录')
    docx_parser.add_argument('--include-cover', action='store_true', help='包含封面')
    
    args = parser.parse_args()
    
    if args.command == 'ppt':
        generate_ppt(args.knowledge, args.output, args.topic, args.template)
    elif args.command == 'docx':
        generate_docx(args.knowledge, args.output, args.topic, 
                     args.include_toc, args.include_cover)
    else:
        parser.print_help()


if __name__ == '__main__':
    main()
